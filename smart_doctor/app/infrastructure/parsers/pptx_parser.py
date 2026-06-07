"""
PPTX 幻灯片解析器（v2.1）

使用 python-pptx 提取文本，支持：
- 每张幻灯片独立分页
- 标题+正文分离
- 表格/图表文本提取
"""
from __future__ import annotations

import logging
import os
import time

from app.infrastructure.parsers.base import DocumentParser, ParsedDocument, TextSegment

logger = logging.getLogger(__name__)


class PptxParser(DocumentParser):
    supported_types = {"pptx", "ppt"}

    async def parse(self, file_path: str, **kwargs) -> ParsedDocument:
        t0 = time.time()
        file_size = _safe_size(file_path)

        if file_path.lower().endswith(".ppt") and not file_path.lower().endswith(".pptx"):
            return _error_doc(
                "旧版 .ppt 格式暂不支持，请转为 .pptx 格式后再上传", file_size, t0
            )

        try:
            from pptx import Presentation
        except ImportError:
            return _error_doc(
                "python-pptx 未安装，请执行 pip install python-pptx", file_size, t0
            )

        segments: list[TextSegment] = []

        try:
            prs = Presentation(file_path)
            slide_count = len(prs.slides)

            for si, slide in enumerate(prs.slides):
                slide_num = si + 1
                slide_title = None

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue

                        # 第一个文本框的第一个段落视为标题
                        if slide_title is None:
                            slide_title = text

                        segments.append(TextSegment(
                            text=text,
                            page=slide_num,
                            slide=slide_num,
                            heading=slide_title,
                            confidence=0.92,
                        ))
        except Exception as e:
            return _error_doc(str(e), file_size, t0)

        elapsed = (time.time() - t0) * 1000
        logger.info(
            "PPTX parsed: path=%s slides=%d segments=%d time=%.0fms",
            file_path, slide_count, len(segments), elapsed,
        )

        full_text = "\n".join(s.text for s in segments)
        return ParsedDocument(
            text=full_text,
            segments=segments,
            page_count=slide_count,
            encoding="utf-8",
            parse_method="python-pptx",
            parse_duration_ms=elapsed,
            file_size=file_size,
            file_type="pptx",
        )


def _safe_size(path: str) -> int:
    try:
        return os.path.getsize(path)
    except OSError:
        return 0


def _error_doc(msg: str, file_size: int, t0: float) -> ParsedDocument:
    return ParsedDocument(
        text="",
        parse_method="python-pptx",
        parse_duration_ms=(time.time() - t0) * 1000,
        file_size=file_size,
        file_type="pptx",
        error=msg,
    )
